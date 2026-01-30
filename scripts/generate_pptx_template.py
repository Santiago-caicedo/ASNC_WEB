"""
Script para generar la plantilla de PowerPoint de ASNC
Ejecutar: python scripts/generate_pptx_template.py

Requiere: pip install python-pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Colores ASNC
NAVY = RGBColor(0x1B, 0x2A, 0x41)      # #1B2A41
GOLD = RGBColor(0xF4, 0xC3, 0x43)      # #F4C343
WHITE = RGBColor(0xFF, 0xFF, 0xFF)     # #FFFFFF
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)  # #F5F7FA
GRAY = RGBColor(0x64, 0x74, 0x8B)      # #64748B


def add_footer(slide, prs):
    """Agrega footer con info de ASNC"""
    footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(9), Inches(0.3)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Asociación Nuclear Colombiana | www.asncol.com"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER


def create_title_slide(prs):
    """Slide 1: Portada"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Fondo navy en la parte superior (60%)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(4.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # Barra dorada decorativa
    gold_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(4.5), Inches(13.33), Inches(0.15)
    )
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()

    # Título principal
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.5), Inches(11.83), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TÍTULO DE LA PRESENTACIÓN"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(3.0), Inches(11.83), Inches(0.8)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Subtítulo o descripción breve"
    p.font.size = Pt(24)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Logo placeholder (texto indicativo)
    logo_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(5.0), Inches(2.5), Inches(1)
    )
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[Logo ASNC]"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Fecha
    date_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(6.2), Inches(11.83), Inches(0.5)
    )
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Fecha de presentación"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER


def create_agenda_slide(prs):
    """Slide 2: Agenda"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.3)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.35), Inches(11.83), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AGENDA"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Barra dorada
    gold_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.15), Inches(1.5), Inches(0.08)
    )
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()

    # Items de agenda
    agenda_items = [
        "01  |  Primer punto de la agenda",
        "02  |  Segundo punto de la agenda",
        "03  |  Tercer punto de la agenda",
        "04  |  Cuarto punto de la agenda",
        "05  |  Quinto punto de la agenda"
    ]

    y_position = 2.0
    for item in agenda_items:
        item_box = slide.shapes.add_textbox(
            Inches(1), Inches(y_position), Inches(10), Inches(0.6)
        )
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = NAVY
        y_position += 0.8

    add_footer(slide, prs)


def create_content_slide(prs):
    """Slide 3: Contenido general"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.3)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.35), Inches(11.83), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TÍTULO DE LA SECCIÓN"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Barra dorada
    gold_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.15), Inches(1.5), Inches(0.08)
    )
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()

    # Contenido
    content_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Contenido principal"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(20)

    # Bullet points
    bullets = [
        "Primer punto importante a destacar en esta sección.",
        "Segundo punto con información relevante.",
        "Tercer punto que complementa la información.",
        "Cuarto punto para cerrar la idea principal."
    ]

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = "•  " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.space_before = Pt(12)

    add_footer(slide, prs)


def create_two_columns_slide(prs):
    """Slide 4: Dos columnas"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.3)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.35), Inches(11.83), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "COMPARACIÓN / DOS COLUMNAS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Barra dorada
    gold_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.15), Inches(1.5), Inches(0.08)
    )
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()

    # Columna izquierda - caja
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT_BG
    left_box.line.color.rgb = GOLD
    left_box.line.width = Pt(2)

    # Título columna izquierda
    left_title = slide.shapes.add_textbox(
        Inches(1), Inches(2.0), Inches(5), Inches(0.6)
    )
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Columna Izquierda"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Contenido columna izquierda
    left_content = slide.shapes.add_textbox(
        Inches(1), Inches(2.7), Inches(5), Inches(3)
    )
    tf = left_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "•  Punto uno\n•  Punto dos\n•  Punto tres"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY

    # Columna derecha - caja
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = NAVY
    right_box.line.fill.background()

    # Título columna derecha
    right_title = slide.shapes.add_textbox(
        Inches(7.25), Inches(2.0), Inches(5), Inches(0.6)
    )
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Columna Derecha"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Contenido columna derecha
    right_content = slide.shapes.add_textbox(
        Inches(7.25), Inches(2.7), Inches(5), Inches(3)
    )
    tf = right_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "•  Punto uno\n•  Punto dos\n•  Punto tres"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

    add_footer(slide, prs)


def create_contact_slide(prs):
    """Slide 5: Contacto"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Fondo navy completo
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Barra dorada superior
    gold_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.15)
    )
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()

    # Logo placeholder
    logo_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(1.0), Inches(2.5), Inches(1)
    )
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[Logo ASNC]"
    p.font.size = Pt(16)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.2), Inches(11.83), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "¡GRACIAS!"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(3.2), Inches(11.83), Inches(0.6)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Asociación Nuclear Colombiana"
    p.font.size = Pt(24)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Línea separadora
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.0), Inches(4.33), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Info de contacto
    contact_info = [
        "www.asncol.com",
        "info@asncol.com",
        "Bucaramanga, Santander, Colombia"
    ]

    y_pos = 4.3
    for info in contact_info:
        info_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(y_pos), Inches(11.83), Inches(0.5)
        )
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        p.text = info
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        y_pos += 0.45

    # Redes sociales
    social_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.8), Inches(11.83), Inches(0.5)
    )
    tf = social_box.text_frame
    p = tf.paragraphs[0]
    p.text = "LinkedIn  •  Instagram  •  Facebook  •  X"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER


def main():
    """Genera la presentación ASNC"""
    # Crear presentación 16:9
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Crear slides
    create_title_slide(prs)
    create_agenda_slide(prs)
    create_content_slide(prs)
    create_two_columns_slide(prs)
    create_contact_slide(prs)

    # Guardar
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, '..', 'static', 'downloads', 'ASNC_Plantilla_Presentacion.pptx')

    # Crear directorio si no existe
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs.save(output_path)
    print(f"✓ Plantilla generada: {output_path}")
    print("\nSlides creados:")
    print("  1. Portada")
    print("  2. Agenda")
    print("  3. Contenido")
    print("  4. Dos Columnas")
    print("  5. Contacto")
    print("\nNota: Reemplaza [Logo ASNC] con el logo real en PowerPoint")


if __name__ == "__main__":
    main()
